from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from odf.opendocument import load as load_od_docs
from odf.text import P as odf_paragraph
from odf import teletype
from odf.draw import Frame
import logging
import os
import hashlib


class StandardEmbedder:
    def __init__(
        self, store, chunk_size=1000, chunk_overlap=200, loglevel=logging.ERROR
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        logging.basicConfig(level=loglevel, handlers=[logging.StreamHandler()])
        self.store = store
        self.logger = logging.getLogger(__name__)

    def extract_text_from_file(self, file_path):
        text_blob = ""
        _, file_extension = os.path.splitext(file_path)
        file_extension = file_extension.lower()
        try:
            if file_extension == ".pdf":
                with open(file_path, "rb") as pdfFile:
                    reader = PdfReader(pdfFile)
                    for page in reader.pages:
                        text_blob += page.extract_text()
            elif file_extension == ".docx":
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text_blob += para.text + "\n"
            elif file_extension == ".pptx":
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_blob += shape.text + "\n"
            elif file_extension == ".odt":
                doc = load_od_docs(file_path)
                for paragraph in doc.getElementsByType(odf_paragraph):
                    text_blob += teletype.extractText(paragraph) + "\n"
            elif file_extension == ".odp":
                doc = load_od_docs(file_path)
                for frame in doc.getElementsByType(Frame):
                    for paragraph in frame.getElementsByType(odf_paragraph):
                        text_blob += teletype.extractText(paragraph) + "\n"
            elif file_extension == ".txt":
                with open(file_path, "r") as txtfile:
                    text_blob += txtfile.read()
            else:
                raise ValueError("Unsupported file format")
        except Exception as ex:
            print(f"An error occurred: {ex}")

        return text_blob

    def train_from_document(self, doc_path: str, doc_id=None):
        text_blob = ""
        if doc_id is None:
            doc_id = hashlib.md5(doc_path.encode("utf-8")).hexdigest()

        try:
            text_blob = self.extract_text_from_file(doc_path)
        except Exception as ex:
            print(ex)
            self.logger.warning(f"Failed to parse: {doc_path}. Error: ", ex)

        if text_blob != "":
            for i, chunk in enumerate(self.chunk_text(text_blob)):
                hashed = f"{doc_path}_doc_{i}".encode("utf-8")
                doc_id = hashlib.md5(hashed).hexdigest()
                self.store.add_document(text=chunk, idx=doc_id)
        else:
            self.logger.warning("All documents failed. Embedding data is empty.")

    def train_from_text(self, text: str, doc_id=None):
        if doc_id is None:
            raise (
                "The vector store needs a doc_id to identify this content. Please pass in 'doc_id', which can be an integer or uuid."
            )

        for chunk in self.chunk_text(text):
            self.store.add_document(text=chunk, idx=doc_id)

    def chunk_text(self, text, buffer_length=1500):
        parts = text.split("\n")
        buffer = ""
        has_flushed = False

        for p in parts:
            if len(buffer) + len(p) >= buffer_length:
                has_flushed = True
                yield buffer.strip("\n")
            if has_flushed:
                has_flushed = False
                buffer = ""
            buffer += "\n" + p

        yield buffer.strip("\n")
